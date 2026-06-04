# -*- coding: utf-8 -*-
import sys, io
sys.stdout = io.TextIOWrapper(sys.stdout.buffer, encoding='utf-8')
from docx import Document
from pptx import Presentation
import openpyxl

base = r"C:\Users\raulh\Desktop\Claude\PlastHub\Context"

def line(c='='): print(c*70)

def read_docx(path, title):
    line(); print("DOCX:", title); line()
    d = Document(path)
    for p in d.paragraphs:
        t = p.text.strip()
        if t:
            style = p.style.name if p.style else ''
            tag = f"[{style}] " if style and style not in ('Normal','Body Text') else ""
            print(tag + t)
    for ti, tbl in enumerate(d.tables):
        print(f"\n--- TABLE {ti+1} ---")
        for row in tbl.rows:
            cells = [c.text.strip().replace('\n',' ') for c in row.cells]
            print(" | ".join(cells))

def read_pptx(path, title):
    line(); print("PPTX:", title); line()
    prs = Presentation(path)
    for si, slide in enumerate(prs.slides):
        print(f"\n===== SLIDE {si+1} =====")
        for shape in slide.shapes:
            if shape.has_text_frame:
                for para in shape.text_frame.paragraphs:
                    t = "".join(r.text for r in para.runs).strip()
                    if t: print("  " + t)
            if shape.has_table:
                tb = shape.table
                print("  [TABLE]")
                for row in tb.rows:
                    print("   " + " | ".join(c.text.strip().replace('\n',' ') for c in row.cells))
        if slide.has_notes_slide and slide.notes_slide.notes_text_frame.text.strip():
            print("  [NOTES]", slide.notes_slide.notes_text_frame.text.strip())

def read_xlsx(path, title):
    line(); print("XLSX:", title); line()
    wb = openpyxl.load_workbook(path, data_only=True)
    for ws in wb.worksheets:
        print(f"\n===== SHEET: {ws.title}  (dims {ws.dimensions}) =====")
        rows = list(ws.iter_rows(values_only=True))
        for r in rows:
            if any(v is not None and str(v).strip()!="" for v in r):
                cells = ["" if v is None else str(v) for v in r]
                # trim trailing empties
                while cells and cells[-1]=="": cells.pop()
                print(" | ".join(cells))

read_xlsx(base + r"\PlasHub_Project_Profile_v2.xlsx", "PlasHub_Project_Profile_v2")
read_pptx(base + r"\Requerimientos_Almacen_Peru.pptx", "Requerimientos_Almacen_Peru")
read_docx(base + r"\Amaru-I-Project-Charter.docx", "Amaru-I-Project-Charter (OUTDATED NUMBERS)")
read_docx(base + "\\Factory “Amaru I” - Proyecto Piloto.docx", "Factory Amaru I - Proyecto Piloto")
